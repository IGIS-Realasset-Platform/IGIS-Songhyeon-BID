import json
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_rgb(rgb_str):
    if not rgb_str:
        return RGBColor(255, 255, 255)
    nums = re.findall(r'\d+', rgb_str)
    if len(nums) >= 3:
        try:
            r = min(max(0, int(nums[0])), 255)
            g = min(max(0, int(nums[1])), 255)
            b = min(max(0, int(nums[2])), 255)
            return RGBColor(r, g, b)
        except Exception as e:
            print(f"Error parsing color '{rgb_str}': {e}")
            return RGBColor(255, 255, 255)
    return RGBColor(255, 255, 255)

def parse_pixel_size(size_str):
    if not size_str:
        return 16.0
    nums = re.findall(r'[\d\.]+', size_str)
    if nums:
        return float(nums[0])
    return 16.0

def is_transparent(color_str):
    if not color_str or color_str == 'transparent':
        return True
    nums = re.findall(r'[\d\.]+', color_str)
    if len(nums) == 4 and float(nums[3]) == 0.0:
        return True
    return False

def group_texts_into_rows(texts, y_threshold=15):
    """
    Groups a list of text objects into horizontal rows.
    If the Y coordinate difference between texts is within y_threshold, they are placed in the same row.
    Inside each row, texts are sorted by X coordinate.
    """
    if not texts:
        return []
    
    # Sort initially by Y
    sorted_texts = sorted(texts, key=lambda t: t['rect']['y'])
    
    rows = []
    current_row = []
    
    for text in sorted_texts:
        if not current_row:
            current_row.append(text)
        else:
            # Compare with the average Y of the current row
            avg_y = sum(item['rect']['y'] for item in current_row) / len(current_row)
            if abs(text['rect']['y'] - avg_y) < y_threshold:
                current_row.append(text)
            else:
                # Sort current row by X before closing it
                current_row.sort(key=lambda t: t['rect']['x'])
                rows.append(current_row)
                current_row = [text]
                
    if current_row:
        current_row.sort(key=lambda t: t['rect']['x'])
        rows.append(current_row)
        
    # Sort rows by their average Y
    rows.sort(key=lambda r: sum(item['rect']['y'] for item in r) / len(r))
    return rows

def build_presentation(json_path, output_path):
    print(f"Loading layout data from {json_path}...")
    if not os.path.exists(json_path):
        print(f"Error: {json_path} does not exist.")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    print(f"Creating PPTX presentation: {output_path}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    scale_x = 13.333 / 1920.0
    scale_y = 7.5 / 1080.0
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_idx, slide_data in enumerate(slides_data):
        print(f"Building slide {slide_data['slideIndex']}...")
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Background color
        bg_color_str = slide_data.get('bgColor', 'rgb(255, 255, 255)')
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = parse_rgb(bg_color_str)
        
        cards = slide_data.get('cards', [])
        texts = slide_data.get('texts', [])
        
        # 2. Draw Cards
        for card in cards:
            rect = card['rect']
            left = Inches(rect['x'] * scale_x)
            top = Inches(rect['y'] * scale_y)
            width = Inches(rect['w'] * scale_x)
            height = Inches(rect['h'] * scale_y)
            
            border_color = parse_rgb(card.get('borderColor', ''))
            border_width_str = card.get('borderWidth', '')
            border_width_px = parse_pixel_size(border_width_str)
            bg_color_str = card.get('bgColor', '')
            
            # CSS border shorthand parsing (detect top-only borders for dividers)
            border_parts = border_width_str.split()
            is_top_line_only = False
            if len(border_parts) >= 2:
                top_px = parse_pixel_size(border_parts[0])
                right_px = parse_pixel_size(border_parts[1])
                bottom_px = parse_pixel_size(border_parts[2]) if len(border_parts) >= 3 else right_px
                left_px = parse_pixel_size(border_parts[3]) if len(border_parts) == 4 else right_px
                
                # If only border-top is present and background is transparent
                if top_px > 0 and right_px == 0 and bottom_px == 0 and left_px == 0:
                    is_top_line_only = True
            
            if is_top_line_only and is_transparent(bg_color_str):
                # Draw a thin horizontal divider line (using a 1pt height solid shape)
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.0))
                shape.fill.solid()
                shape.fill.fore_color.rgb = border_color
                shape.line.fill.background()
            else:
                # Normal card box
                border_radius_str = card.get('borderRadius', '')
                has_radius = any(px in border_radius_str for px in ['px', 'rem', '%']) and '0px' not in border_radius_str
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if has_radius else MSO_SHAPE.RECTANGLE
                
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                
                if is_transparent(bg_color_str):
                    shape.fill.background()
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = parse_rgb(bg_color_str)
                
                if border_width_px > 0:
                    shape.line.color.rgb = border_color
                    shape.line.width = Pt(max(1.0, border_width_px * 0.75))
                else:
                    shape.line.fill.background()
                
        # 3. Assign texts to cards geometrically based on center point (with nested card priority)
        card_grouped_texts = {i: [] for i in range(len(cards))}
        unassigned_texts = []
        
        for text in texts:
            t_rect = text['rect']
            t_cx = t_rect['x'] + t_rect['w'] / 2
            t_cy = t_rect['y'] + t_rect['h'] / 2
            
            assigned_c_idx = -1
            min_area = float('inf')
            
            for c_idx, card in enumerate(cards):
                c_rect = card['rect']
                # Check if text center is inside card
                if (c_rect['x'] <= t_cx <= c_rect['x'] + c_rect['w'] and 
                    c_rect['y'] <= t_cy <= c_rect['y'] + c_rect['h']):
                    area = c_rect['w'] * c_rect['h']
                    # Keep the card with the smallest area (most specific child card)
                    if area < min_area:
                        min_area = area
                        assigned_c_idx = c_idx
                        
            if assigned_c_idx != -1:
                card_grouped_texts[assigned_c_idx].append(text)
            else:
                unassigned_texts.append(text)
                
        # 4. Render Grouped Card Texts
        for c_idx, grouped_texts in card_grouped_texts.items():
            if not grouped_texts:
                continue
            card_rect = cards[c_idx]['rect']
            
            # Card padding configuration
            padding_x = 10 # px (reduced padding to prevent wrap issues)
            padding_y = 10 # px
            left = Inches((card_rect['x'] + padding_x) * scale_x)
            top = Inches((card_rect['y'] + padding_y) * scale_y)
            # Add extra horizontal safety space to prevent text wrap on right edge
            width = Inches((card_rect['w'] - 2 * padding_x) * scale_x) + Inches(0.4)
            height = Inches((card_rect['h'] - 2 * padding_y) * scale_y) + Inches(0.2)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            # Group card texts into horizontal rows
            text_rows = group_texts_into_rows(grouped_texts, y_threshold=15)
            
            prev_row_bottom_y = card_rect['y'] + padding_y
            
            for r_idx, row in enumerate(text_rows):
                p = tf.paragraphs[0] if r_idx == 0 else tf.add_paragraph()
                
                # Space before calculation
                current_row_y = min(item['rect']['y'] for item in row)
                gap = current_row_y - prev_row_bottom_y
                if gap > 4 and r_idx > 0:
                    p.space_before = Pt(gap * 0.75)
                
                # Apply text alignment of the first item in the row
                text_align = row[0].get('textAlign', 'left')
                if text_align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif text_align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
                # Add runs for each item in the row
                for item_idx, text_item in enumerate(row):
                    # If not first item in row, add spacing based on X gap
                    if item_idx > 0:
                        prev_item = row[item_idx - 1]
                        gap_x = text_item['rect']['x'] - (prev_item['rect']['x'] + prev_item['rect']['w'])
                        if gap_x > 0:
                            fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                            space_char_width = fs_px * 0.28
                            num_spaces = max(1, int(gap_x / space_char_width))
                            run_space = p.add_run()
                            run_space.text = " " * num_spaces
                            run_space.font.size = Pt(fs_px * 0.75)
                            
                    run = p.add_run()
                    run.text = text_item['text']
                    
                    font_family = text_item.get('fontFamily', '').lower()
                    is_kr_text = any(ord(char) > 127 for char in text_item['text'])
                    
                    if is_kr_text:
                        run.font.name = "Pretendard"
                    elif 'sanomat' in font_family:
                        run.font.name = "Sanomat Wp"
                    else:
                        run.font.name = "Arial"
                        
                    fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                    run.font.size = Pt(fs_px * 0.75)
                    
                    font_weight = text_item.get('fontWeight', '')
                    if font_weight == 'bold' or (font_weight.isdigit() and int(font_weight) >= 600):
                        run.font.bold = True
                        
                    run.font.color.rgb = parse_rgb(text_item.get('color', 'rgb(0,0,0)'))
                
                prev_row_bottom_y = max(item['rect']['y'] + item['rect']['h'] for item in row)

        # 5. Render Unassigned Grouped Column Texts (Titles, Subtitles, etc.)
        # Group unassigned texts by vertical columns
        unassigned_texts.sort(key=lambda t: t['rect']['y'])
        text_groups = []
        for text in unassigned_texts:
            t_rect = text['rect']
            t_cx = t_rect['x'] + t_rect['w'] / 2
            
            placed = False
            for group in text_groups:
                avg_cx = sum(item['rect']['x'] + item['rect']['w']/2 for item in group) / len(group)
                if abs(t_cx - avg_cx) < 120:
                    group.append(text)
                    placed = True
                    break
            if not placed:
                text_groups.append([text])
                
        for group in text_groups:
            if not group:
                continue
            
            # Calculate original bounding box of the group
            g_left_px = min(item['rect']['x'] for item in group)
            g_top_px = min(item['rect']['y'] for item in group)
            g_right_px = max(item['rect']['x'] + item['rect']['w'] for item in group)
            g_bottom_px = max(item['rect']['y'] + item['rect']['h'] for item in group)
            
            g_w = g_right_px - g_left_px
            
            # Check if group contains centered text or header elements to apply massive width buffer
            has_center = any(item.get('textAlign', 'left') == 'center' for item in group)
            has_header = any(item.get('tagName', '') in ['H1', 'H2', 'H3'] for item in group)
            
            if has_center or has_header or g_w > 500:
                # Centered widescreen layout textbox
                g_w_buffered = 1600.0
                g_left_buffered = 960.0 - 800.0
            else:
                g_w_buffered = g_w * 1.40 # 40% buffer
                g_cx = g_left_px + g_w / 2
                g_left_buffered = g_cx - g_w_buffered / 2
            
            # Clamp buffer bounds to slide dimensions
            g_left_buffered = max(40.0, g_left_buffered)
            g_right_buffered = min(1920.0 - 40.0, g_left_buffered + g_w_buffered)
            g_w_buffered = g_right_buffered - g_left_buffered
            
            left = Inches(g_left_buffered * scale_x)
            top = Inches(g_top_px * scale_y)
            width = Inches(g_w_buffered * scale_x)
            height = Inches((g_bottom_px - g_top_px) * scale_y) + Inches(0.4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            # Group into horizontal rows
            group_rows = group_texts_into_rows(group, y_threshold=15)
            
            prev_row_bottom_y = g_top_px
            
            for r_idx, row in enumerate(group_rows):
                p = tf.paragraphs[0] if r_idx == 0 else tf.add_paragraph()
                
                # Space before calculation
                current_row_y = min(item['rect']['y'] for item in row)
                gap = current_row_y - prev_row_bottom_y
                if gap > 4 and r_idx > 0:
                    p.space_before = Pt(gap * 0.75)
                
                # Apply text alignment of the first item in the row
                text_align = row[0].get('textAlign', 'left')
                if text_align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif text_align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
                # Add runs for each item in the row
                for item_idx, text_item in enumerate(row):
                    # Add X spacing if not first item
                    if item_idx > 0:
                        prev_item = row[item_idx - 1]
                        gap_x = text_item['rect']['x'] - (prev_item['rect']['x'] + prev_item['rect']['w'])
                        if gap_x > 0:
                            fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                            space_char_width = fs_px * 0.28
                            num_spaces = max(1, int(gap_x / space_char_width))
                            run_space = p.add_run()
                            run_space.text = " " * num_spaces
                            run_space.font.size = Pt(fs_px * 0.75)
                            
                    run = p.add_run()
                    run.text = text_item['text']
                    
                    font_family = text_item.get('fontFamily', '').lower()
                    is_kr_text = any(ord(char) > 127 for char in text_item['text'])
                    
                    if is_kr_text:
                        run.font.name = "Pretendard"
                    elif 'sanomat' in font_family:
                        run.font.name = "Sanomat Wp"
                    else:
                        run.font.name = "Arial"
                        
                    fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                    run.font.size = Pt(fs_px * 0.75)
                    
                    font_weight = text_item.get('fontWeight', '')
                    if font_weight == 'bold' or (font_weight.isdigit() and int(font_weight) >= 600):
                        run.font.bold = True
                        
                    run.font.color.rgb = parse_rgb(text_item.get('color', 'rgb(0,0,0)'))
                
                prev_row_bottom_y = max(item['rect']['y'] + item['rect']['h'] for item in row)
                
    prs.save(output_path)
    print(f"Generated presentation successfully: {output_path}")

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))
    build_presentation(os.path.join(base_dir, 'layout_data_kr.json'), os.path.join(base_dir, 'IOTA_Strategy_KR.pptx'))
    print("Done building presentations.")
