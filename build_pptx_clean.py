import json
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_rgb(rgb_str):
    if not rgb_str:
        return RGBColor(29, 29, 31)
    nums = re.findall(r'\d+', rgb_str)
    if len(nums) >= 3:
        try:
            r = min(max(0, int(nums[0])), 255)
            g = min(max(0, int(nums[1])), 255)
            b = min(max(0, int(nums[2])), 255)
            return RGBColor(r, g, b)
        except Exception as e:
            return RGBColor(29, 29, 31)
    return RGBColor(29, 29, 31)

def parse_pixel_size(size_str):
    if not size_str:
        return 16.0
    nums = re.findall(r'[\d\.]+', size_str)
    if nums:
        return float(nums[0])
    return 16.0

def is_transparent(color_str):
    if not color_str or color_str == 'transparent':
        return True
    nums = re.findall(r'[\d\.]+', color_str)
    if len(nums) == 4 and float(nums[3]) == 0.0:
        return True
    return False

def is_dark_color(rgb_color):
    return (rgb_color[0] + rgb_color[1] + rgb_color[2]) / 3.0 < 120.0

def build_presentation(json_path, output_path):
    print(f"Loading layout data from {json_path}...")
    if not os.path.exists(json_path):
        print(f"Error: {json_path} does not exist.")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    print(f"Creating Clean Template PPTX presentation: {output_path}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    for slide_data in slides_data:
        slide_idx = slide_data['slideIndex']
        print(f"Building slide {slide_idx}...")
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Background color detection & Dark Mode setup
        bg_color_str = slide_data.get('bgColor', 'rgb(253, 253, 253)')
        bg_rgb = parse_rgb(bg_color_str)
        
        # Strict standardized backgrounds for clean corporate presentation look
        is_dark = is_dark_color(bg_rgb)
        final_bg_color = RGBColor(29, 29, 31) if is_dark else RGBColor(253, 253, 253)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = final_bg_color
        
        # Text default colors
        text_color_main = RGBColor(255, 255, 255) if is_dark else RGBColor(29, 29, 31)
        text_color_sub = RGBColor(160, 160, 160) if is_dark else RGBColor(110, 110, 110)
        accent_blue = RGBColor(96, 165, 250) if is_dark else RGBColor(30, 58, 138)
        
        cards = slide_data.get('cards', [])
        texts = slide_data.get('texts', [])
        
        # 2. Extract Titles & Subtitles (Header Area)
        header_texts = []
        content_texts = []
        for text in texts:
            if text['rect']['y'] < 380:
                header_texts.append(text)
            else:
                content_texts.append(text)
                
        # Sort header texts vertically to identify Theme/Subtitle and Main Title
        header_texts.sort(key=lambda t: t['rect']['y'])
        
        theme_text = ""
        title_text = ""
        
        if len(header_texts) >= 2:
            t1_size = parse_pixel_size(header_texts[0].get('fontSize', '16px'))
            t2_size = parse_pixel_size(header_texts[1].get('fontSize', '16px'))
            if t1_size < t2_size:
                theme_text = header_texts[0]['text']
                title_text = header_texts[1]['text']
            else:
                title_text = header_texts[0]['text']
                theme_text = header_texts[1]['text']
        elif len(header_texts) == 1:
            title_text = header_texts[0]['text']
            
        # 3. Render Standardized Slide Header
        if theme_text or title_text:
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.5))
            tf = header_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            is_first = True
            if theme_text:
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = theme_text
                run.font.name = "Pretendard"
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = text_color_sub
                is_first = False
                
            if title_text:
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = title_text
                run.font.name = "Pretendard"
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = text_color_main
                
        # 4. Process Cards & Dividers
        clean_cards = []
        dividers = []
        
        for card in cards:
            rect = card['rect']
            bg_color_str = card.get('bgColor', '')
            border_width_str = card.get('borderWidth', '')
            
            if rect['w'] >= 1700 and rect['h'] >= 900:
                continue # Skip outer slide container
                
            border_parts = border_width_str.split()
            is_top_line_only = False
            if len(border_parts) >= 2:
                top_px = parse_pixel_size(border_parts[0])
                right_px = parse_pixel_size(border_parts[1])
                bottom_px = parse_pixel_size(border_parts[2]) if len(border_parts) >= 3 else right_px
                left_px = parse_pixel_size(border_parts[3]) if len(border_parts) == 4 else right_px
                if top_px > 0 and right_px == 0 and bottom_px == 0 and left_px == 0:
                    is_top_line_only = True
                    
            if is_top_line_only and is_transparent(bg_color_str):
                dividers.append(card)
            elif rect['h'] < 15:
                dividers.append(card)
            else:
                clean_cards.append(card)
                
        # 5. Determine Layout Mode & Render Content
        M = len(clean_cards)
        content_texts.sort(key=lambda t: t['rect']['y'])
        
        if M == 0:
            # TEXT-ONLY LIST LAYOUT
            content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(10.333), Inches(4.3))
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            # Group text nodes into rows
            rows = []
            current_row = []
            for text in content_texts:
                if not current_row:
                    current_row.append(text)
                else:
                    avg_y = sum(item['rect']['y'] for item in current_row) / len(current_row)
                    if abs(text['rect']['y'] - avg_y) < 15:
                        current_row.append(text)
                    else:
                        current_row.sort(key=lambda t: t['rect']['x'])
                        rows.append(current_row)
                        current_row = [text]
            if current_row:
                current_row.sort(key=lambda t: t['rect']['x'])
                rows.append(current_row)
                
            for r_idx, row in enumerate(rows):
                p = tf.paragraphs[0] if r_idx == 0 else tf.add_paragraph()
                p.space_after = Pt(14)
                
                align = row[0].get('textAlign', 'left')
                p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
                
                for item_idx, item in enumerate(row):
                    if item_idx > 0:
                        prev_item = row[item_idx - 1]
                        gap_x = item['rect']['x'] - (prev_item['rect']['x'] + prev_item['rect']['w'])
                        if gap_x > 0:
                            num_spaces = max(2, int(gap_x / 12))
                            run_space = p.add_run()
                            run_space.text = " " * num_spaces
                            
                    run = p.add_run()
                    run.text = item['text']
                    run.font.name = "Pretendard"
                    
                    fs_px = parse_pixel_size(item.get('fontSize', '18px'))
                    run.font.size = Pt(max(15, min(24, fs_px * 0.9)))
                    
                    if item.get('fontWeight', '') == 'bold' or (item.get('fontWeight', '').isdigit() and int(item.get('fontWeight', '')) >= 600):
                        run.font.bold = True
                    run.font.color.rgb = text_color_main
                    
        else:
            # CARD/GRID LAYOUT
            card_texts = {i: [] for i in range(M)}
            for text in content_texts:
                t_cx = text['rect']['x'] + text['rect']['w'] / 2
                t_cy = text['rect']['y'] + text['rect']['h'] / 2
                
                closest_idx = 0
                min_dist = float('inf')
                for c_idx, card in enumerate(clean_cards):
                    c_rect = card['rect']
                    c_cx = c_rect['x'] + c_rect['w'] / 2
                    c_cy = c_rect['y'] + c_rect['h'] / 2
                    dist = (t_cx - c_cx)**2 + (t_cy - c_cy)**2
                    if dist < min_dist:
                        min_dist = dist
                        closest_idx = c_idx
                card_texts[closest_idx].append(text)
                
            # Grid structure calculations
            if M == 1:
                cols, rows = 1, 1
            elif M == 2:
                cols, rows = 2, 1
            elif M == 3:
                cols, rows = 3, 1
            elif M == 4:
                cols, rows = 4, 1
            elif M <= 6:
                cols, rows = 3, 2
            else:
                cols = 4
                rows = (M + 3) // 4
                
            c_x_start = 1.0
            c_y_start = 2.4
            c_w_total = 11.333
            c_h_total = 4.4
            
            gap_x = 0.4
            gap_y = 0.3
            
            card_w = (c_w_total - (cols - 1) * gap_x) / cols
            card_h = (c_h_total - (rows - 1) * gap_y) / rows
            
            # Sort cards physically to position them left-to-right, top-to-bottom
            orig_cards_sorted = sorted(range(M), key=lambda idx: (clean_cards[idx]['rect']['y'], clean_cards[idx]['rect']['x']))
            
            for grid_pos_idx, orig_card_idx in enumerate(orig_cards_sorted):
                col_idx = grid_pos_idx % cols
                row_idx = grid_pos_idx // cols
                
                left = c_x_start + col_idx * (card_w + gap_x)
                top = c_y_start + row_idx * (card_h + gap_y)
                
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
                
                orig_card = clean_cards[orig_card_idx]
                orig_card_bg = parse_rgb(orig_card.get('bgColor', ''))
                
                if is_transparent(orig_card.get('bgColor', '')) or orig_card_bg == bg_rgb:
                    # Glassmorphism/Transparent card styling
                    shape.fill.background()
                    shape.line.color.rgb = accent_blue
                    shape.line.width = Pt(1.5)
                else:
                    shape.fill.solid()
                    if is_dark:
                        shape.fill.fore_color.rgb = RGBColor(45, 45, 48)
                        shape.line.fill.background()
                    else:
                        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
                        shape.line.color.rgb = RGBColor(226, 232, 240)
                        shape.line.width = Pt(1.0)
                        
                tx_padding_x = 0.15
                tx_padding_y = 0.15
                txBox = slide.shapes.add_textbox(
                    Inches(left + tx_padding_x), 
                    Inches(top + tx_padding_y), 
                    Inches(card_w - 2 * tx_padding_x), 
                    Inches(card_h - 2 * tx_padding_y)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
                
                t_list = card_texts[orig_card_idx]
                t_list.sort(key=lambda t: t['rect']['y'])
                
                card_rows = []
                curr_card_row = []
                for item in t_list:
                    if not curr_card_row:
                        curr_card_row.append(item)
                    else:
                        avg_y = sum(x['rect']['y'] for x in curr_card_row) / len(curr_card_row)
                        if abs(item['rect']['y'] - avg_y) < 15:
                            curr_card_row.append(item)
                        else:
                            curr_card_row.sort(key=lambda x: x['rect']['x'])
                            card_rows.append(curr_card_row)
                            curr_card_row = [item]
                if curr_card_row:
                    curr_card_row.sort(key=lambda x: x['rect']['x'])
                    card_rows.append(curr_card_row)
                    
                for cr_idx, c_row in enumerate(card_rows):
                    p = tf.paragraphs[0] if cr_idx == 0 else tf.add_paragraph()
                    p.space_after = Pt(6)
                    
                    align = c_row[0].get('textAlign', 'left')
                    p.alignment = PP_ALIGN.CENTER if align == 'center' else PP_ALIGN.LEFT
                    
                    for item_idx, item in enumerate(c_row):
                        if item_idx > 0:
                            prev_item = c_row[item_idx - 1]
                            gap_x = item['rect']['x'] - (prev_item['rect']['x'] + prev_item['rect']['w'])
                            if gap_x > 0:
                                run_space = p.add_run()
                                run_space.text = " " * max(1, int(gap_x / 10))
                                
                        run = p.add_run()
                        run.text = item['text']
                        run.font.name = "Pretendard"
                        
                        fs_px = parse_pixel_size(item.get('fontSize', '15px'))
                        run.font.size = Pt(max(12, min(18, fs_px * 0.85)))
                        
                        if item.get('fontWeight', '') == 'bold' or (item.get('fontWeight', '').isdigit() and int(item.get('fontWeight', '')) >= 600):
                            run.font.bold = True
                            
                        # Format index numbers with accent color
                        is_index = len(item['text']) <= 3 and item['text'].strip().replace('.', '').isdigit()
                        if is_index:
                            run.font.color.rgb = accent_blue
                            run.font.bold = True
                        else:
                            run.font.color.rgb = text_color_main
                            
            # Render Dividers elegantly at standardized heights
            for div in dividers:
                rect = div['rect']
                y_scale = rect['y'] / 1080.0
                top_inch = y_scale * 7.5
                
                # Draw standard divider line
                slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    Inches(1.0), 
                    Inches(top_inch), 
                    Inches(11.333), 
                    Pt(1.0)
                )
                
    prs.save(output_path)
    print(f"Generated clean standard template presentation successfully: {output_path}")

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))
    build_presentation(os.path.join(base_dir, 'layout_data_kr.json'), os.path.join(base_dir, 'IOTA_Strategy_Clean_KR.pptx'))
    print("Done building presentations.")
