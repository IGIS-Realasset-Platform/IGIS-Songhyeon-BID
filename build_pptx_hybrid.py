import json
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_rgb(color_str):
    if not color_str or color_str == 'transparent':
        return RGBColor(30, 58, 138) # Default to brand blue: 1e3a8a
        
    color_str = color_str.strip().lower()
    
    # 1. Hex color #rrggbb or #rgb
    if color_str.startswith('#'):
        hex_val = color_str.lstrip('#')
        if len(hex_val) == 3:
            hex_val = ''.join([c*2 for c in hex_val])
        if len(hex_val) == 6:
            try:
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                # CRITICAL: Detect neon green (e.g. #00ff00) and correct it to brand blue
                if g > 180 and r < 50 and b < 50:
                    return RGBColor(30, 58, 138)
                return RGBColor(r, g, b)
            except Exception:
                return RGBColor(30, 58, 138)
                
    # 2. rgb(r, g, b) or rgba(r, g, b, a)
    if color_str.startswith('rgb'):
        nums = re.findall(r'[\d\.]+', color_str)
        if len(nums) >= 3:
            try:
                r = min(max(0, int(float(nums[0]))), 255)
                g = min(max(0, int(float(nums[1]))), 255)
                b = min(max(0, int(float(nums[2]))), 255)
                # CRITICAL: Detect and correct neon green bug
                if g > 180 and r < 50 and b < 50:
                    return RGBColor(30, 58, 138)
                return RGBColor(r, g, b)
            except Exception:
                return RGBColor(30, 58, 138)
                
    # 3. Fallback for any leftover unparsed colors
    return RGBColor(30, 58, 138)

def parse_pixel_size(size_str):
    if not size_str:
        return 16.0
    nums = re.findall(r'[\d\.]+', size_str)
    if nums:
        return float(nums[0])
    return 16.0

def group_texts_into_rows(texts, y_threshold=15):
    if not texts:
        return []
    
    # Sort initially by Y
    sorted_texts = sorted(texts, key=lambda t: t['rect']['y'])
    
    rows = []
    current_row = []
    
    for text in sorted_texts:
        if not current_row:
            current_row.append(text)
        else:
            avg_y = sum(item['rect']['y'] for item in current_row) / len(current_row)
            if abs(text['rect']['y'] - avg_y) < y_threshold:
                current_row.append(text)
            else:
                current_row.sort(key=lambda t: t['rect']['x'])
                rows.append(current_row)
                current_row = [text]
                
    if current_row:
        current_row.sort(key=lambda t: t['rect']['x'])
        rows.append(current_row)
        
    rows.sort(key=lambda r: sum(item['rect']['y'] for item in r) / len(r))
    return rows

def build_presentation(json_path, output_path):
    print(f"Loading layout data from {json_path}...")
    if not os.path.exists(json_path):
        print(f"Error: {json_path} does not exist.")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    print(f"Creating Hybrid Presentation: {output_path}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    scale_x = 13.333 / 1920.0
    scale_y = 7.5 / 1080.0
    
    blank_layout = prs.slide_layouts[6]
    base_dir = os.path.dirname(os.path.abspath(__file__))
    
    for slide_data in slides_data:
        slide_idx = slide_data['slideIndex']
        print(f"Building slide {slide_idx}...")
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Insert transparent text-free design template as background picture
        bg_image_path = os.path.join(base_dir, 'img', f'slide_bg_{slide_idx}.png')
        if os.path.exists(bg_image_path):
            slide.shapes.add_picture(bg_image_path, 0, 0, Inches(13.333), Inches(7.5))
            print(f"Inserted background template: {bg_image_path}")
        else:
            print(f"Warning: Background template not found: {bg_image_path}")
            # Fallback to solid color background
            bg_color_str = slide_data.get('bgColor', 'rgb(255, 255, 255)')
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = parse_rgb(bg_color_str)
            
        cards = slide_data.get('cards', [])
        texts = slide_data.get('texts', [])
        
        # Note: We do NOT draw cards on top because they are already perfectly rendered 
        # in the background picture. This completely bypasses PPTX styling constraints.
        
        # 2. Render Texts - Split every text element into individual textboxes (No Paragraph/Card Grouping)
        for text_item in texts:
            rect = text_item['rect']
            
            # Identify header text elements to widen their bounds (prevents premature wrapping)
            is_header_text = rect['y'] < 380
            
            if is_header_text:
                # Center header textbox with 1600px width buffer
                t_w_px = 1600.0
                left = Inches((960.0 - 800.0) * scale_x)
                width = Inches(t_w_px * scale_x)
            else:
                # Widen textboxes dynamically based on text length & alignment to prevent wrapping
                orig_x = rect['x']
                orig_w = rect['w']
                
                # Apply 1.8x buffer for long description text to prevent ugly wraps
                text_len = len(text_item.get('text', ''))
                factor = 1.80 if text_len > 15 else 1.45
                new_w = orig_w * factor
                
                text_align = text_item.get('textAlign', 'left')
                if text_align == 'center':
                    cx = orig_x + orig_w / 2
                    left_px = cx - new_w / 2
                elif text_align == 'right':
                    right_edge = orig_x + orig_w
                    left_px = right_edge - new_w
                else:
                    left_px = orig_x
                
                left = Inches(left_px * scale_x)
                width = Inches((new_w + 40.0) * scale_x)
                
            top = Inches(rect['y'] * scale_y)
            # Add minor height padding
            height = Inches(rect['h'] * scale_y) + Inches(0.1)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            p = tf.paragraphs[0]
            
            # Apply alignment
            text_align = text_item.get('textAlign', 'left')
            if text_align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
            # Apply Paragraph-level default font for editing safety
            try:
                p.font.name = "Pretendard Variable"
                fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                p.font.size = Pt(fs_px * 0.75)
            except Exception:
                pass
                
            run = p.add_run()
            run.text = text_item['text']
            
            # Strict Font Requirement: Pretendard Variable
            run.font.name = "Pretendard Variable"
            
            # Font size (px to pt)
            fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
            run.font.size = Pt(fs_px * 0.75)
            
            # Font weight
            font_weight = text_item.get('fontWeight', '')
            if font_weight == 'bold' or (font_weight.isdigit() and int(font_weight) >= 600):
                run.font.bold = True
                
            # Text color
            run.font.color.rgb = parse_rgb(text_item.get('color', 'rgb(0,0,0)'))
                
    prs.save(output_path)
    print(f"Generated hybrid presentation successfully: {output_path}")

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))
    build_presentation(os.path.join(base_dir, 'layout_data_kr.json'), os.path.join(base_dir, 'IOTA_Strategy_Hybrid_KR.pptx'))
    print("Done building presentations.")
