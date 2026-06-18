import json
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_rgb(color_str):
    if not color_str or color_str == 'transparent':
        return RGBColor(30, 58, 138) # Default to brand blue: 1e3a8a
        
    color_str = color_str.strip().lower()
    
    # 1. Hex color #rrggbb or #rgb
    if color_str.startswith('#'):
        hex_val = color_str.lstrip('#')
        if len(hex_val) == 3:
            hex_val = ''.join([c*2 for c in hex_val])
        if len(hex_val) == 6:
            try:
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                # CRITICAL: Detect neon green (e.g. #00ff00) and correct it to brand blue
                if g > 180 and r < 50 and b < 50:
                    return RGBColor(30, 58, 138)
                return RGBColor(r, g, b)
            except Exception:
                return RGBColor(30, 58, 138)
                
    # 2. rgb(r, g, b) or rgba(r, g, b, a)
    if color_str.startswith('rgb'):
        nums = re.findall(r'[\d\.]+', color_str)
        if len(nums) >= 3:
            try:
                r = min(max(0, int(float(nums[0]))), 255)
                g = min(max(0, int(float(nums[1]))), 255)
                b = min(max(0, int(float(nums[2]))), 255)
                # CRITICAL: Detect and correct neon green bug (0, 255, 0 or similar green fallbacks)
                if g > 180 and r < 50 and b < 50:
                    return RGBColor(30, 58, 138)
                return RGBColor(r, g, b)
            except Exception:
                return RGBColor(30, 58, 138)
                
    # 3. Fallback for any leftover unparsed colors
    return RGBColor(30, 58, 138)

def parse_pixel_size(size_str):
    if not size_str:
        return 16.0
    nums = re.findall(r'[\d\.]+', size_str)
    if nums:
        return float(nums[0])
    return 16.0

def is_transparent(color_str):
    if not color_str or color_str == 'transparent':
        return True
    nums = re.findall(r'[\d\.]+', color_str)
    if len(nums) == 4 and float(nums[3]) == 0.0:
        return True
    return False

def build_presentation(json_path, output_path):
    print(f"Loading layout data from {json_path}...")
    if not os.path.exists(json_path):
        print(f"Error: {json_path} does not exist.")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    print(f"Creating Pure Native PPTX presentation (Split Textboxes & Brand Blue): {output_path}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    scale_x = 13.333 / 1920.0
    scale_y = 7.5 / 1080.0
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    for slide_data in slides_data:
        slide_idx = slide_data['slideIndex']
        print(f"Building slide {slide_idx}...")
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Set background color (Pure native color)
        bg_color_str = slide_data.get('bgColor', 'rgb(255, 255, 255)')
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = parse_rgb(bg_color_str)
        
        cards = slide_data.get('cards', [])
        texts = slide_data.get('texts', [])
        
        # 2. Draw Cards (Rectangles/Rounded Rectangles) - Pure Native Shapes
        for card in cards:
            rect = card['rect']
            left = Inches(rect['x'] * scale_x)
            top = Inches(rect['y'] * scale_y)
            width = Inches(rect['w'] * scale_x)
            height = Inches(rect['h'] * scale_y)
            
            border_color = parse_rgb(card.get('borderColor', ''))
            border_width_str = card.get('borderWidth', '')
            border_width_px = parse_pixel_size(border_width_str)
            bg_color_str = card.get('bgColor', '')
            
            # CSS border shorthand parsing (detect top-only borders for dividers)
            border_parts = border_width_str.split()
            is_top_line_only = False
            if len(border_parts) >= 2:
                top_px = parse_pixel_size(border_parts[0])
                right_px = parse_pixel_size(border_parts[1])
                bottom_px = parse_pixel_size(border_parts[2]) if len(border_parts) >= 3 else right_px
                left_px = parse_pixel_size(border_parts[3]) if len(border_parts) == 4 else right_px
                
                # If only border-top is present and background is transparent
                if top_px > 0 and right_px == 0 and bottom_px == 0 and left_px == 0:
                    is_top_line_only = True
            
            if is_top_line_only and is_transparent(bg_color_str):
                # Draw a thin horizontal divider line (using a 1pt height solid shape)
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.0))
                shape.fill.solid()
                shape.fill.fore_color.rgb = border_color
                shape.line.fill.background()
            else:
                # Normal card box
                border_radius_str = card.get('borderRadius', '')
                has_radius = any(px in border_radius_str for px in ['px', 'rem', '%']) and '0px' not in border_radius_str
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if has_radius else MSO_SHAPE.RECTANGLE
                
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                
                # Apply fill color
                if is_transparent(bg_color_str):
                    shape.fill.background()
                else:
                    shape.fill.solid()
                    card_bg = parse_rgb(bg_color_str)
                    shape.fill.fore_color.rgb = card_bg
                
                # Apply border/line formatting
                if border_width_px > 0:
                    shape.line.color.rgb = border_color
                    shape.line.width = Pt(max(1.0, border_width_px * 0.75))
                else:
                    shape.line.fill.background()
                
        # 3. Render Texts - Split every text element into individual textboxes (No Paragraph Merging)
        for text_item in texts:
            rect = text_item['rect']
            
            # Identify header text elements to widen their bounds (prevents premature wrapping)
            is_header_text = rect['y'] < 380
            
            if is_header_text:
                # Center header textbox with 1600px width buffer
                t_w_px = 1600.0
                left = Inches((960.0 - 800.0) * scale_x)
                width = Inches(t_w_px * scale_x)
            else:
                # Widen textboxes dynamically based on text length & alignment to prevent wrapping
                orig_x = rect['x']
                orig_w = rect['w']
                
                # Apply 1.8x buffer for long description text to prevent ugly wraps
                text_len = len(text_item.get('text', ''))
                factor = 1.80 if text_len > 15 else 1.45
                new_w = orig_w * factor
                
                text_align = text_item.get('textAlign', 'left')
                if text_align == 'center':
                    cx = orig_x + orig_w / 2
                    left_px = cx - new_w / 2
                elif text_align == 'right':
                    right_edge = orig_x + orig_w
                    left_px = right_edge - new_w
                else:
                    left_px = orig_x
                
                left = Inches(left_px * scale_x)
                width = Inches((new_w + 40.0) * scale_x)
                
            top = Inches(rect['y'] * scale_y)
            # Add minor height padding
            height = Inches(rect['h'] * scale_y) + Inches(0.1)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
            
            p = tf.paragraphs[0]
            
            # Apply alignment
            text_align = text_item.get('textAlign', 'left')
            if text_align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
            # Apply Paragraph-level default font for editing safety
            try:
                p.font.name = "Pretendard Variable"
                fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
                p.font.size = Pt(fs_px * 0.75)
            except Exception:
                pass
                
            run = p.add_run()
            run.text = text_item['text']
            
            # Strict Font Requirement: Pretendard Variable
            run.font.name = "Pretendard Variable"
            
            # Font size (px to pt)
            fs_px = parse_pixel_size(text_item.get('fontSize', '16px'))
            run.font.size = Pt(fs_px * 0.75)
            
            # Font weight
            font_weight = text_item.get('fontWeight', '')
            if font_weight == 'bold' or (font_weight.isdigit() and int(font_weight) >= 600):
                run.font.bold = True
                
            # Text color
            run.font.color.rgb = parse_rgb(text_item.get('color', 'rgb(0,0,0)'))
            
    prs.save(output_path)
    print(f"Generated pure native split-textbox presentation successfully: {output_path}")

if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.abspath(__file__))
    build_presentation(os.path.join(base_dir, 'layout_data_kr.json'), os.path.join(base_dir, 'IOTA_Strategy_Native_KR.pptx'))
    print("Done building presentations.")
