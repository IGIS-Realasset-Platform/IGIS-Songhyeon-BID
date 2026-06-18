import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def get_slide_order(main_layout_path):
    print(f"Reading slide sequence from {main_layout_path}...")
    with open(main_layout_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Extract the slides array
    match = re.search(r'const\s+slides\s*=\s*React\.useMemo\(\(\)\s*=>\s*\[([\s\S]*?)\]', content)
    if not match:
        # Fallback to general list search
        match = re.search(r'\[\s*(<Section1\s*/>[\s\S]*?)\]', content)
        
    if match:
        slide_array_str = match.group(1)
        slide_names = re.findall(r'<(\w+)\s*/?>', slide_array_str)
        print(f"Found {len(slide_names)} slides in active presentation.")
        return slide_names
    else:
        print("Warning: Could not parse slides array from MainLayout.jsx. Listing file-based sections...")
        # Fallback to sorted names of files in components
        components_dir = os.path.dirname(main_layout_path)
        files = os.listdir(components_dir)
        section_files = [f.split('.')[0] for f in files if f.startswith('Section') and f.endswith('.jsx')]
        # Sort numerically if possible, otherwise alphabetically
        def try_int(s):
            m = re.search(r'\d+', s)
            return int(m.group(0)) if m else s
        section_files.sort(key=try_int)
        return section_files

def clean_html_tags(text):
    text = text.strip()
    # Remove fragments, br, strong, span, strong tags
    text = re.sub(r'</?.*?>', ' ', text)
    # Remove HTML entities
    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
    # Unquote if the whole string is single or double quoted
    if (text.startswith("'") and text.endswith("'")) or (text.startswith('"') and text.endswith('"')):
        text = text[1:-1]
    # Remove extra spaces
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def parse_ternary_expression(expr_body):
    # Splits {lang === 'kr' ? KR_PART : EN_PART} by top-level ':'
    depth = 0
    split_idx = -1
    for i, char in enumerate(expr_body):
        if char in ('{', '('):
            depth += 1
        elif char in ('}', ')'):
            depth -= 1
        elif char == ':' and depth == 0:
            split_idx = i
            break
            
    if split_idx != -1:
        kr_part = expr_body[:split_idx].strip()
        en_part = expr_body[split_idx+1:].strip()
        
        # Clean parentheses
        if kr_part.startswith('(') and kr_part.endswith(')'):
            kr_part = kr_part[1:-1].strip()
        if en_part.startswith('(') and en_part.endswith(')'):
            en_part = en_part[1:-1].strip()
            
        return kr_part, en_part
    return expr_body, expr_body

def extract_texts_from_jsx(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return [], []
        
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Check if this is a dark slide heuristically (look for dark bg classes)
    is_dark = False
    if 'bg-[#1d1d1f]' in content or 'bg-gray-900' in content or 'bg-black' in content:
        is_dark = True
        
    # Remove style tags
    content = re.sub(r'<style>[\s\S]*?</style>', '', content)
    
    # Match the JSX return statement
    return_match = re.search(r'return\s*\(\s*(<[\s\S]*?>)\s*\)\s*;?\s*\n?\s*\}', content)
    if not return_match:
        return_match = re.search(r'return\s*\(([\s\S]*)\)\s*;?\s*\}', content)
        
    jsx_body = return_match.group(1) if return_match else content
    
    kr_texts = []
    en_texts = []
    
    i = 0
    n = len(jsx_body)
    
    in_tag = False
    brace_depth_in_tag = 0
    current_text = []
    
    while i < n:
        char = jsx_body[i]
        
        if char == '<' and not in_tag:
            if current_text:
                txt = "".join(current_text).strip()
                if txt:
                    cleaned = clean_html_tags(txt)
                    if cleaned:
                        kr_texts.append(cleaned)
                        en_texts.append(cleaned)
                current_text = []
            in_tag = True
            brace_depth_in_tag = 0
            i += 1
            continue
            
        elif in_tag and char == '{':
            brace_depth_in_tag += 1
            i += 1
            continue
            
        elif in_tag and char == '}':
            brace_depth_in_tag = max(0, brace_depth_in_tag - 1)
            i += 1
            continue
            
        elif char == '>' and in_tag and brace_depth_in_tag == 0:
            in_tag = False
            i += 1
            continue
            
        if in_tag:
            i += 1
        else:
            if char == '{':
                if current_text:
                    txt = "".join(current_text).strip()
                    if txt:
                        cleaned = clean_html_tags(txt)
                        if cleaned:
                            kr_texts.append(cleaned)
                            en_texts.append(cleaned)
                    current_text = []
                
                brace_depth = 1
                j = i + 1
                expr = []
                while j < n and brace_depth > 0:
                    c = jsx_body[j]
                    if c == '{':
                        brace_depth += 1
                    elif c == '}':
                        brace_depth -= 1
                    if brace_depth > 0:
                        expr.append(c)
                    j += 1
                
                expr_str = "".join(expr).strip()
                # Check for language translation
                if 'lang === \'kr\'' in expr_str or 'lang === "kr"' in expr_str or 'lang===\'kr\'' in expr_str:
                    cond_pattern = r'lang\s*===\s*[\'"]kr[\'"]\s*\?\s*'
                    match = re.search(cond_pattern, expr_str)
                    if match:
                        expr_body = expr_str[match.end():].strip()
                        kr_part, en_part = parse_ternary_expression(expr_body)
                        cleaned_kr = clean_html_tags(kr_part)
                        cleaned_en = clean_html_tags(en_part)
                        if cleaned_kr:
                            kr_texts.append(cleaned_kr)
                        if cleaned_en:
                            en_texts.append(cleaned_en)
                else:
                    # Not a translation expression, skip state conditions
                    pass
                i = j
            else:
                current_text.append(char)
                i += 1
                
    if current_text:
        txt = "".join(current_text).strip()
        if txt:
            cleaned = clean_html_tags(txt)
            if cleaned:
                kr_texts.append(cleaned)
                en_texts.append(cleaned)
                
    # Final cleanup list
    def final_cleanup(lst):
        res = []
        for x in lst:
            x = re.sub(r'</?\s*.*?>', ' ', x)
            x = re.sub(r'\s+', ' ', x).strip()
            if not x or x in ('true', 'false', 'null', 'undefined'):
                continue
            if x.startswith('const') or x.startswith('return') or x.startswith('import') or x.startswith('export'):
                continue
            res.append(x)
        return res
        
    return final_cleanup(kr_texts), final_cleanup(en_texts), is_dark

def create_presentation(slide_names, components_dir, output_path, is_korean=True):
    print(f"Creating presentation at {output_path}...")
    prs = Presentation()
    # Set widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout (layout index 6 is typically blank)
    blank_layout = prs.slide_layouts[6]
    
    for idx, name in enumerate(slide_names):
        file_path = os.path.join(components_dir, f"{name}.jsx")
        kr_texts, en_texts, is_dark = extract_texts_from_jsx(file_path)
        
        texts = kr_texts if is_korean else en_texts
        
        # Add slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Determine background color
        background = slide.background
        fill = background.fill
        fill.solid()
        if is_dark:
            fill.fore_color.rgb = RGBColor(0x1D, 0x1D, 0x1F) # Apple dark background
            text_color = RGBColor(0xFF, 0xFF, 0xFF)
            sub_color = RGBColor(0xA1, 0xA1, 0xA1)
        else:
            fill.fore_color.rgb = RGBColor(0xFD, 0xFD, 0xFD) # Apple light background
            text_color = RGBColor(0x1D, 0x1D, 0x1F)
            sub_color = RGBColor(0x88, 0x88, 0x88)
            
        if not texts:
            # Empty slide fallback
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{name} (No text / Image Slide)"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = sub_color
            p.alignment = PP_ALIGN.CENTER
            continue
            
        # Parse content hierarchy:
        # Typically: Theme/Chapter, Main Title, and then list items.
        theme = ""
        title = ""
        bullets = []
        
        # Try to identify chapter/theme
        if len(texts) >= 1:
            first = texts[0]
            if first.startswith("Chapter") or first.startswith("Part") or first.startswith("[") or len(first) < 30:
                theme = first
                if len(texts) >= 2:
                    title = texts[1]
                    bullets = texts[2:]
                else:
                    title = ""
                    bullets = []
            else:
                title = first
                bullets = texts[1:]
                
        # 1. Slide Title Text Box (Theme + Main Title)
        # Position at the top: Left 1.0", Top 1.0", Width 11.333", Height 2.2"
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(2.2))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        
        if theme:
            p_theme = tf_title.paragraphs[0]
            p_theme.text = theme
            p_theme.font.name = "Malgun Gothic" if is_korean else "Arial"
            p_theme.font.size = Pt(16)
            p_theme.font.bold = True
            p_theme.font.color.rgb = sub_color
            
            if title:
                p_title = tf_title.add_paragraph()
                p_title.text = title
                p_title.font.name = "Malgun Gothic" if is_korean else "Arial"
                p_title.font.size = Pt(32)
                p_title.font.bold = True
                p_title.font.color.rgb = text_color
                p_title.space_before = Pt(12)
        elif title:
            p_title = tf_title.paragraphs[0]
            p_title.text = title
            p_title.font.name = "Malgun Gothic" if is_korean else "Arial"
            p_title.font.size = Pt(32)
            p_title.font.bold = True
            p_title.font.color.rgb = text_color
            
        # 2. Slide Body Text Box (Bullet points / paragraphs)
        # Position below title: Left 1.5", Top 3.2", Width 10.333", Height 3.5"
        if bullets:
            body_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.333), Inches(3.5))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            tf_body.margin_left = tf_body.margin_right = tf_body.margin_top = tf_body.margin_bottom = 0
            
            for b_idx, bullet in enumerate(bullets):
                # Clean up bullet characters if they are parsed literally (like '▪')
                bullet = bullet.strip().lstrip('▪').lstrip('-').strip()
                if not bullet:
                    continue
                
                p_bullet = tf_body.paragraphs[0] if b_idx == 0 and len(tf_body.paragraphs[0].text) == 0 else tf_body.add_paragraph()
                p_bullet.text = "•  " + bullet
                p_bullet.font.name = "Malgun Gothic" if is_korean else "Arial"
                p_bullet.font.size = Pt(18)
                p_bullet.font.color.rgb = text_color
                p_bullet.space_after = Pt(12)
                p_bullet.line_spacing = 1.3
                
        # If it's a simple slide with only title and no body, vertically center it
        if not bullets and not theme and title:
            # Move the title box to the center
            title_box.top = Inches(2.8)
            title_box.height = Inches(2.0)
            tf_title.paragraphs[0].alignment = PP_ALIGN.CENTER
        elif not bullets and theme and title:
            title_box.top = Inches(2.5)
            title_box.height = Inches(2.5)
            tf_title.paragraphs[0].alignment = PP_ALIGN.CENTER
            if len(tf_title.paragraphs) > 1:
                tf_title.paragraphs[1].alignment = PP_ALIGN.CENTER
                
    prs.save(output_path)
    print(f"Presentation saved successfully with {len(prs.slides)} slides.")

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    main_layout_path = os.path.join(base_dir, 'src', 'components', 'MainLayout.jsx')
    components_dir = os.path.join(base_dir, 'src', 'components')
    
    slide_names = get_slide_order(main_layout_path)
    
    if not slide_names:
        print("Error: Slide names list is empty.")
        return
        
    output_kr = os.path.join(base_dir, 'IOTA_Strategy_KR.pptx')
    output_en = os.path.join(base_dir, 'IOTA_Strategy_EN.pptx')
    
    # Generate Korean presentation
    create_presentation(slide_names, components_dir, output_kr, is_korean=True)
    
    # Generate English presentation
    create_presentation(slide_names, components_dir, output_en, is_korean=False)
    
    print("\n--- ALL EXPORTS COMPLETED ---")
    print(f"Korean: {output_kr}")
    print(f"English: {output_en}")

if __name__ == "__main__":
    main()
